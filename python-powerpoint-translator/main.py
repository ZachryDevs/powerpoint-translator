from flask import Flask, request, send_file, render_template
from pptx import Presentation
from googletrans import Translator
import os
import tempfile

app = Flask(__name__)

@app.route('/', methods=['GET'])
def index():
    return render_template('upload_form.html')

def translate_text(text, target_lang):
    translator = Translator()
    translated = translator.translate(text, src='en', dest=target_lang)
    return translated.text

def translate_presentation(input_path, output_path, target_lang, font_adjustment='auto'):
    prs = Presentation(input_path)
    i = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                if str(original_text) != '' and original_text is not None:
                    for char in str(original_text):
                        if char.isalnum():
                            i = True
                            break
                    if i:
                        translated_text = translate_text(original_text, target_lang)
                        
                        # Handle font size adjustment based on user preference
                        if font_adjustment == 'auto' and hasattr(shape, "text_frame"):
                            length_ratio = len(translated_text) / len(original_text)
                            if length_ratio > 1.2:  # Text is longer
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if hasattr(run, "font"):
                                            current_size = run.font.size
                                            if current_size:
                                                new_size = int(current_size / length_ratio)
                                                run.font.size = new_size
                            elif length_ratio < 0.8:  # Text is shorter
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if hasattr(run, "font"):
                                            current_size = run.font.size
                                            if current_size:
                                                new_size = int(current_size / length_ratio)
                                                run.font.size = new_size
                        
                        shape.text = translated_text
                        i = False

    prs.save(output_path)


@app.route('/translate', methods=['POST'])
def translate_pptx_and_return():
    uploaded_file = request.files['file']
    target_lang = request.form.get('target_language', 'es')  # Default to Spanish if not specified
    font_adjustment = request.form.get('font_adjustment', 'auto')  # Default to automatic if not specified

    if uploaded_file.filename != '':
        temp_dir = tempfile.mkdtemp()
        uploaded_path = os.path.join(temp_dir, uploaded_file.filename)
        uploaded_file.save(uploaded_path)

        ppt_name = os.path.splitext(uploaded_file.filename)[0]
        output_path = os.path.join(temp_dir, f"{ppt_name}-translated({target_lang}).pptx")

        translate_presentation(uploaded_path, output_path, target_lang, font_adjustment)

        return send_file(output_path, as_attachment=True, download_name=f"{ppt_name}-translated({target_lang}).pptx")

    return "No file uploaded"

if __name__ == "__main__":
    app.run(debug=True)

# from pptx import Presentation
# from googletrans import Translator
#
#
# def translate_to_spanish(text):
#     translator = Translator()
#     translated = translator.translate(text, src='en', dest='es')
#     return translated.text
#
#
# def translate_presentation(input_path, output_path):
#     prs = Presentation(input_path)
#     i = False
#
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 original_text = shape.text
#                 if str(original_text) != '' and original_text is not None:
#                     for char in str(original_text):
#                         if char.isalnum():
#                             i = True
#                             break
#                     if i:
#                         translated_text = translate_to_spanish(original_text)
#                         shape.text = translated_text
#                         i = False
#
#     prs.save(output_path)
#
#
# if __name__ == "__main__":
#     ppt_name = "8_14_Origin_of_the_Universe"
#     input_pptx = ppt_name + ".pptx"
#     output_pptx = ppt_name + "-translated(es).pptx"
#
#     translate_presentation(input_pptx, output_pptx)
